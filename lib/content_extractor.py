"""按扩展名提取文件文本预览（客户端独立副本，不依赖后端仓库）。"""

from __future__ import annotations

import os
from pathlib import Path

MAX_CHARS = int(os.getenv("DOCMIND_CONTENT_PREVIEW_CHARS", "2000"))

SUPPORTED = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".md": "txt",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
}


def extract_preview(file_path: str, *, max_chars: int | None = None) -> str:
    limit = max_chars if max_chars is not None else MAX_CHARS
    path = Path(file_path)
    if not path.exists():
        return ""

    ext = path.suffix.lower()
    kind = SUPPORTED.get(ext)
    if not kind:
        return f"[不支持的格式: {ext}]"

    try:
        if kind == "txt":
            text = path.read_text(encoding="utf-8", errors="ignore")
        elif kind == "pdf":
            text = _extract_pdf(path)
        elif kind == "docx":
            text = _extract_docx(path)
        elif kind == "xlsx":
            text = _extract_xlsx(path)
        elif kind == "pptx":
            text = _extract_pptx(path)
        elif kind == "image":
            text = _extract_image_ocr(path)
        else:
            text = ""
    except Exception as exc:
        return f"[提取失败: {exc}]"

    return text[:limit]


def _extract_pdf(path: Path) -> str:
    import pdfplumber

    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages[:5]:
            parts.append(page.extract_text() or "")
    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    rows = []
    for sheet in wb.worksheets[:2]:
        for row in sheet.iter_rows(max_row=20, values_only=True):
            rows.append(" ".join(str(c) for c in row if c is not None))
    return "\n".join(rows)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    texts = []
    for slide in prs.slides[:10]:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def _extract_image_ocr(path: Path) -> str:
    import pytesseract
    from PIL import Image

    return pytesseract.image_to_string(Image.open(path), lang="chi_sim+eng")
